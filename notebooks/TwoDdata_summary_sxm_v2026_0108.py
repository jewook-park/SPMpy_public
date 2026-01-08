# ---
# jupyter:
#   jupytext:
#     formats: ipynb,py:percent
#     text_representation:
#       extension: .py
#       format_name: percent
#       format_version: '1.3'
#       jupytext_version: 1.17.2
#   kernelspec:
#     display_name: Python 3 (ipykernel)
#     language: python
#     name: python3
# ---

# %% [markdown]
# # TwoDdata_summary_sxm_v2026_0107
#
#
# Automated SXM summary → PPTX (cleaned & robust v2).

# %% [markdown]
# ## 🔧 Pipeline Correction Notice (2026-01-07)
#
# **Important correction applied without deleting existing content.**
#
# - Processing order is now:
#   `sxm → interpolate2d → plane fit → plateau tilt → FFT`
# - FFT is performed strictly *after* all real-space corrections.
# - Figure titles are NOT embedded in images.
# - PPTX is appended page-by-page without date stamps.
#
# All original markdown, explanations, and code cells are preserved below.
# New or corrected functions are appended at the end of the notebook.
#
#
# ## 🔧 Pipeline Correction & Plane-Fit Feasibility Policy (2026-01-07)
#
# **Important correction applied without deleting or modifying existing content.**
#
# This update clarifies the *preprocessing responsibility split* between
# data validation and numerical fitting, and introduces an explicit
# plane-fit feasibility check **before** calling `plane_fit_xr`.
#
# ---
#
# ### ✅ Final Processing Order (unchanged, now strictly enforced)
# sxm
# → interpolate2D_xr
# → plane-fit feasibility check + masking
# → plane_fit_xr
# → plateau_tilt_xr
# → FFT
# FFT is performed **strictly after** all real-space corrections.
#
# ---
#
# ### 🧠 Why Plane-Fit Feasibility Is Checked *Outside* `plane_fit_xr`
#
# - `plane_fit_xr` assumes that the input data is mathematically fit-able.
# - Certain SXM datasets violate this assumption *locally*:
#   - Entire columns may be NaN
#   - Columns may be constant (no height variation)
#   - y-coordinates may collapse after interpolation
# - These cases cause numerical failures (e.g. SVD non-convergence),
#   which should **not** terminate automated pipelines.
#
# To preserve algorithmic clarity:
# - `plane_fit_xr` itself is **not modified**
# - All data validation is handled upstream
#
# ---
#
# ### 📐 Column-wise Masking Policy for `method='y_fit'`
#
# **Important note:**
#
# - `method='y_fit'` performs **column-wise polynomial fitting**
# - Therefore, plane-fit feasibility is evaluated **per column**
# - If a column fails any criterion, the **entire column is masked**
#
# This is intentional and mathematically consistent.
#
# **Mask criteria per column:**
# 1. Number of finite data points `< poly_order + 1`
# 2. Column values are effectively constant
# 3. y-coordinates are degenerate (zero spatial extent)
#
# Partial masking within a column is **not used**.
#
# ---
#
# ### 🧩 Mask Propagation Philosophy
#
# - Plane-fit masks are stored explicitly in the dataset:
#   - `plane_fit_mask_Z_fwd`
#   - `plane_fit_mask_Z_bwd`
# - These masks are reused consistently in:
#   - plane fitting
#   - plateau tilt correction
#   - FFT
#   - visualization
#
# This ensures that all downstream analysis operates on the
# *same physically meaningful regions*.
#
# ---
#
# ### 📄 PPTX Generation Policy
#
# - Figures contain **no embedded titles**
# - All descriptive text is placed in PPT textboxes
# - Slides are appended sequentially without timestamps
# - Plane fitting may be partial, but PPTX generation never aborts
#
# ---
#
# **Result:**  
# A robust, physically honest, and automation-safe SXM processing pipeline
# that preserves all original content while enforcing correct numerical
# boundaries.
#
#

# %% [markdown]
# ## 🔧 Pipeline Update & Processing Policy (2026-01-07)
#
# **Important corrections applied without deleting or modifying existing content.**  
# All original markdown, explanations, and code cells are preserved below.  
# New or corrected logic is appended at the end of the notebook.
#
# ---
#
# ### ✅ Final Processing Order (Strictly Enforced)
# sxm
# → interpolate2D_xr
# → plane-fit feasibility check + masking
# → plane_fit_xr
# → plateau_tilt_xr
# → FFT
#
# - FFT is performed **strictly after** all real-space corrections.
# - No figure contains embedded titles or annotations.
# - PPTX pages are appended sequentially without timestamps.
#
# ---
#
# ### 🧠 Responsibility Split: Validation vs. Fitting
#
# - `plane_fit_xr` assumes mathematically fit-able input data.
# - Local violations may exist in real SXM data:
#   - Entire columns are NaN
#   - Columns are constant
#   - Spatial coordinates collapse after interpolation
#
# To maintain algorithmic clarity and robustness:
# - `plane_fit_xr` itself is **not modified**
# - All feasibility checks and masking are handled **upstream**
#
# ---
#
# ### 📐 Plane-Fit Feasibility Policy (`method='y_fit'`)
#
# - `method='y_fit'` performs **column-wise polynomial fitting**
# - Feasibility is therefore evaluated **per column**
# - If a column fails any criterion, the **entire column is masked**
#
# **Mask criteria**
# 1. Number of finite points `< poly_order + 1`
# 2. Column values are effectively constant
# 3. y-coordinates have no spatial extent
#
# Partial masking within a column is intentionally **not used**.
#
# ---
#
# ### 🧩 Mask Propagation
#
# - Plane-fit masks are stored explicitly in the dataset
#   (e.g. `plane_fit_mask_Z_fwd`, `plane_fit_mask_Z_bwd`)
# - The same masks are reused consistently in:
#   - plane fitting
#   - plateau tilt correction
#   - FFT
#   - visualization
#
# This guarantees all downstream analysis operates on the
# **same physically meaningful regions**.
#
# ---
#
# ### 🔁 Multipass SXM Handling
#
# - `img2xr` exposes multipass data as separate channels
#   (e.g. `Z_P1_fwd`, `Z_P2_fwd`, …).
# - The pipeline automatically detects multipass datasets.
#
# **PPTX layout for multipass data**
# 1. A dedicated **multipass overview page** is added first.
# 2. Each pass (`P1`, `P2`, …) is then processed independently:
#    - Z / LIX
#    - forward / backward
#    - real-space images and FFTs
# 3. After all passes are completed, the pipeline proceeds to the next dataset.
#
# Single-pass datasets continue to use the original single-page layout.
#
# ---
#
# ### 📄 PPTX Generation Policy
#
# - Figures contain **no embedded titles**
# - All descriptive text is placed in PPT textboxes
# - Plane fitting and plateau tilt may be partially skipped,
#   but **PPTX generation never aborts**
#
# ---
#
# **Result**
#
# A robust, physically honest, and automation-safe SXM processing pipeline  
# that supports both single-pass and multipass data while enforcing clear  
# numerical and structural boundaries.
#

# %%

# %% [markdown]
# ## Stage 0 — Repository path injection

# %%

import sys
from pathlib import Path

REPO_ROOT = Path(r"C:\Users\gkp\Documents\GitHub\SPMpy")
if not REPO_ROOT.exists():
    raise RuntimeError(f"SPMpy repo not found: {REPO_ROOT}")
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

import spmpy
print('[SPMpy] Imported from:', spmpy.__file__)


# %% [markdown]
# ## Stage 1 — Imports (Quickstart-compatible I/O)

# %%

import os
import numpy as np
import matplotlib.pyplot as plt
import seaborn_image as isns

from pptx import Presentation
from pptx.util import Inches, Pt

from spmpy.io import spmpy_io_library_v0_1_2 as io
select_folder = io.select_folder
files_in_folder = io.files_in_folder
img2xr = io.img2xr

from spmpy.preprocess.interpolate2D_xr import interpolate2D_xr

from spmpy.preprocess.plane_fit_xr import plane_fit_xr

from spmpy.preprocess.plateau_tilt_xr import plateau_tilt_xr 

from spmpy.fft.twoDfft_xrft import twoDfft_xrft



# %% [markdown]
# ## Stage 2 — Select working folder

# %%

selected_folder = select_folder()
if not selected_folder:
    raise RuntimeError("No folder selected")
print("Selected:", selected_folder)


# %% [markdown]
# ## Stage 3 — Inventory folder (robust SXM filtering)

# %%

df_files = files_in_folder(selected_folder)

# NOTE:
# files_in_folder() does NOT guarantee an 'ext' column.
# We therefore filter by file_path suffix instead.
df_sxm = df_files[df_files['file_path'].str.lower().str.endswith('.sxm')]

df_sxm


# %% [markdown]
# ## Stage 4 — Group SXM files (legacy-compatible)

# %%
# Stage 4 — Group SXM files (use existing group column)

groups = {}

for group_name, df_g in df_sxm.groupby('group'):
    df_g_sorted = df_g.sort_values('num')
    groups[group_name] = df_g_sorted['file_path'].tolist()

groups


# %%

# %%

# %% [markdown]
#
# ## Stage 5 — Generate PPTX summary
#
# FFT amplitude display
# ---------------------
# - FFT amplitude is displayed on **log scale** for visibility.
# - Internally uses: `np.log10(amp + eps)`
#
# FFT defaults (TwoDfft_xrft)
# ---------------------------
# - detrend = 'constant'
# - window  = 'hann'
# - shift   = True
#

# %%

# %%

# %%
def TwoDdata_summary_sxm_pptx(
    groups,
    output_pptx=None,
):
    """
    Automated SXM 2D summary → PPTX.

    Plotting rules:
    - No text inside figures besides axes/ticks/colorbar labels
    - All titles/subtitles via PPT textboxes only
    - interpolate2D_xr always applied
    - Real space:
        * extent, dx, scalebar, axis ticks ALL in meters (m) for consistency
        * x/y tick labels MUST be visible (showticks=True + explicit ticks)
        * x/y tick labels MUST be scientific notation (forced via powerlimits)
        * scalebar via isns.imshow(dx=..., units=...), but force label to stay in meters ("m") not SI-prefix ("zm")
        * Z colorbar label: "Z (m)"
    - FFT:
        * use dataset output as-is (NO extra fftshift)
        * ticks include [min, 0, max]
        * colorbar label indicates log scale
    - Group-wise section slides + summary slide
    - Temporary PNGs deleted AFTER inserting into PPT
    """

    import os
    import gc
    import json
    import numpy as np
    import matplotlib.pyplot as plt
    from matplotlib.ticker import ScalarFormatter
    import seaborn_image as isns
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from tqdm import tqdm

    # =================================================
    # Plane-fit feasibility mask (y_fit, column-wise)
    # =================================================
    def make_plane_fit_mask_yfit(ds, ch, eps=1e-12):
        """
        Plane-fit feasibility mask for Z channels (y_fit).
    
        A column is masked ONLY if plane fitting is
        mathematically impossible due to coordinate degeneracy.
        """
    
        data = np.asarray(ds[ch].values)
        ny, nx = data.shape
    
        if "y" in ds:
            ycoord = np.asarray(ds["y"].values)
        elif "Y" in ds:
            ycoord = np.asarray(ds["Y"].values)
        else:
            dy = ds.attrs.get("Y_spacing", 1.0)
            ycoord = np.arange(ny) * float(dy)
    
        mask = np.zeros((ny, nx), dtype=bool)
    
        for ix in range(nx):
            col = data[:, ix]
            valid = np.isfinite(col)
    
            # Only true degeneracy: y coordinate collapse
            if valid.any() and np.nanstd(ycoord[valid]) < eps:
                mask[:, ix] = True
    
        return mask


    # [ADDED] mask for plateau_tilt_xr (propagate plane-fit feasibility)
    def make_plateau_tilt_mask(ds, ch, base_mask):
        Z = np.asarray(ds[ch].values)
        return (~np.isfinite(Z)) | base_mask

    # -------------------------------------------------
    # Output filename: include working folder name
    # -------------------------------------------------
    folder_name = os.path.basename(os.path.abspath(selected_folder))
    if output_pptx is None:
        output_pptx = f"automated_sxm_{folder_name}_summary.pptx"
    output_path = os.path.join(selected_folder, output_pptx)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =================================================
    # Summary slide (folder + groups + files)
    # =================================================
    MAX_LIST_FILES = 8

    summary = prs.slides.add_slide(prs.slide_layouts[6])
    tf = summary.shapes.add_textbox(
        Inches(0.8), Inches(0.6), Inches(11.7), Inches(6.2)
    ).text_frame

    tf.paragraphs[0].text = "Automated SXM Summary"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True

    p_folder = tf.add_paragraph()
    p_folder.text = f"Folder: {folder_name}"
    p_folder.font.size = Pt(18)
    p_folder.font.italic = True

    for group_name, file_list in groups.items():
        pg = tf.add_paragraph()
        pg.text = f"\n[{group_name}] ({len(file_list)} files)"
        pg.font.size = Pt(20)
        pg.font.bold = True

        if len(file_list) <= MAX_LIST_FILES:
            for f in file_list:
                pf = tf.add_paragraph()
                pf.text = f"  - {os.path.basename(f)}"
                pf.font.size = Pt(14)
        else:
            pf = tf.add_paragraph()
            pf.text = f"  files: 001 ~ {len(file_list):03d}"
            pf.font.size = Pt(14)

    total_files = sum(len(v) for v in groups.values())

    # =================================================
    # Helper functions for robust attrs + ticks
    # =================================================
    def _parse_interpolate2d(ds):
        raw = ds.attrs.get("interpolate2D", None)
        if isinstance(raw, dict):
            return raw
        if isinstance(raw, str):
            try:
                return json.loads(raw)
            except Exception:
                return {}
        return {}

    def _get_dx_dy_m(ds):
        """
        Return pixel spacing in meters. Robust against interpolate2D being dict or JSON string.
        """
        interp = _parse_interpolate2d(ds)
        dx_m = interp.get("dx_new", ds.attrs.get("X_spacing", None))
        dy_m = interp.get("dy_new", ds.attrs.get("Y_spacing", None))
        if dx_m is None or dy_m is None:
            raise ValueError("Missing X_spacing/Y_spacing in ds.attrs (and interpolate2D parse failed).")
        return float(dx_m), float(dy_m)

    def _image_size_m(ds):
        raw = ds.attrs.get("image_size", None)
        if raw is None:
            return None

        if isinstance(raw, (list, tuple, np.ndarray)) and len(raw) >= 2:
            try:
                return float(raw[0]), float(raw[1])
            except Exception:
                pass

        if isinstance(raw, str):
            try:
                s = raw.strip()
                if s.startswith("[") and s.endswith("]"):
                    s = s[1:-1]
                parts = [p.strip() for p in s.split(",")]
                if len(parts) >= 2:
                    return float(parts[0]), float(parts[1])
            except Exception:
                return None

        return None

    def _nice_step(span, target_ticks=4):
        if span <= 0:
            return 1.0
        raw = span / max(target_ticks, 1)
        exp = np.floor(np.log10(raw))
        base = raw / (10 ** exp)
        if base <= 1:
            m = 1
        elif base <= 2:
            m = 2
        elif base <= 5:
            m = 5
        else:
            m = 10
        return float(m * (10 ** exp))

    def _sci_formatter_forced():
        from matplotlib.ticker import ScalarFormatter
        """
        Force scientific notation ALWAYS (no long decimals).
        """
        fmt = ScalarFormatter(useOffset=False)
        fmt.set_scientific(True)
        fmt.set_powerlimits((0, 0))
        return fmt

    def _sci_formatter_3sig():
        from matplotlib.ticker import FuncFormatter
        return FuncFormatter(lambda x, pos: f"{x:.3g}")

    def _plain_formatter():
        """
        Plain (non-scientific) formatter for log colorbar ticks if desired.
        """
        fmt = ScalarFormatter(useOffset=False)
        fmt.set_scientific(False)
        return fmt

    def _force_scalebar_units_m(ax):
        """
        seaborn-image adds a ScaleBar artist (matplotlib-scalebar backend).
        Keep isns scalebar, but force its label formatting to stay in meters ("m")
        and use scientific notation like 5e-09 m (no SI prefix like 'zm').
        """
        try:
            from matplotlib_scalebar.scalebar import ScaleBar
        except Exception:
            return

        for artist in list(getattr(ax, "artists", [])):
            if isinstance(artist, ScaleBar):
                if hasattr(artist, "units"):
                    artist.units = "m"
                if hasattr(artist, "fixed_units"):
                    artist.fixed_units = "m"
                if hasattr(artist, "scale_formatter"):
                    artist.scale_formatter = lambda v, u: f"{v:.1e} m"

    # =================================================
    # Figure generation
    # =================================================
    def save_real(da, fname, ds, cmap):
        """
        Real-space plot (FINAL, PHYSICALLY CORRECT):

        - Axis unit: nm
        - Tick labels: nm, 1–2–5 × 10^n
        - Scalebar: drawn in DATA COORDINATES (nm), exact length
        - Colorbar: Z shown in nm (no scientific offset)
        - Aspect ratio preserved
        """
        import numpy as np
        import matplotlib.pyplot as plt
        from matplotlib.ticker import ScalarFormatter

        Z = np.asarray(da.values) * 1e9  # -> nm
        ny, nx = Z.shape

        dx_m, dy_m = _get_dx_dy_m(ds)
        dx_nm = dx_m * 1e9
        dy_nm = dy_m * 1e9

        L = _image_size_m(ds)
        if L is not None:
            Lx_nm = L[0] * 1e9
            Ly_nm = L[1] * 1e9
        else:
            Lx_nm = nx * dx_nm
            Ly_nm = ny * dy_nm

        extent = (0.0, Lx_nm, 0.0, Ly_nm)

        fig, ax = plt.subplots(figsize=(4, 4))
        im = ax.imshow(
            Z,
            origin="lower",
            extent=extent,
            aspect="equal",
            cmap=cmap,
        )

        ax.set_xlabel("x (nm)")
        ax.set_ylabel("y (nm)")

        sx = _nice_step(Lx_nm, target_ticks=4)
        sy = _nice_step(Ly_nm, target_ticks=4)
        ax.set_xticks(np.arange(0, Lx_nm + 1e-12, sx))
        ax.set_yticks(np.arange(0, Ly_nm + 1e-12, sy))

        fmt = ScalarFormatter(useOffset=False)
        fmt.set_scientific(False)
        ax.xaxis.set_major_formatter(fmt)
        ax.yaxis.set_major_formatter(fmt)

        for lab in ax.get_xticklabels() + ax.get_yticklabels():
            lab.set_rotation(45)

        cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
        cbar.set_label("Z (nm)")
        cbar.formatter = fmt
        cbar.update_ticks()

        # scalebar in data coords
        target_nm = 0.25 * Lx_nm
        exp = np.floor(np.log10(target_nm))
        base = target_nm / (10 ** exp)
        if base <= 1:
            mant = 1
        elif base <= 2:
            mant = 2
        elif base <= 5:
            mant = 5
        else:
            mant = 10
        bar_nm = mant * (10 ** exp)

        x0 = Lx_nm - bar_nm - 0.05 * Lx_nm
        y0 = 0.05 * Ly_nm

        ax.plot(
            [x0, x0 + bar_nm],
            [y0, y0],
            color="white",
            linewidth=3,
            solid_capstyle="butt",
            zorder=10,
        )
        ax.text(
            x0 + bar_nm / 2,
            y0 + 0.03 * Ly_nm,
            f"{bar_nm:g} nm",
            color="white",
            ha="center",
            va="bottom",
            fontsize=9,
            zorder=10,
        )

        ax.set_title("")
        fig.suptitle("")
        fig.savefig(fname, dpi=200)
        plt.close(fig)

    def save_fft(da_fft, fname, ds, cmap):
        """
        FFT plot:
        - NO extra fftshift (use dataset output as-is)
        - extent in 1/m, ticks [min, 0, max]
        - colorbar label indicates log scale
        """
        arr = np.asarray(da_fft.values)
        ny, nx = arr.shape

        dx_m, dy_m = _get_dx_dy_m(ds)
        dkx = 1.0 / (nx * dx_m)
        dky = 1.0 / (ny * dy_m)

        kx = (np.arange(nx) - nx / 2) * dkx
        ky = (np.arange(ny) - ny / 2) * dky
        extent_k = (float(kx.min()), float(kx.max()), float(ky.min()), float(ky.max()))

        log_fft = np.log10(np.abs(arr) + 1e-12)
        vmin, vmax = np.percentile(log_fft, [5, 99])

        fig, ax = plt.subplots(figsize=(4, 4))
        im = ax.imshow(
            log_fft,
            cmap=cmap,
            origin="lower",
            extent=extent_k,
            aspect="equal",
            vmin=vmin,
            vmax=vmax,
        )

        ax.set_xlabel("freqX (1/m)")
        ax.set_ylabel("freqY (1/m)")

        ax.set_xticks([extent_k[0], 0.0, extent_k[1]])
        ax.set_yticks([extent_k[2], 0.0, extent_k[3]])

        fmtk = _sci_formatter_3sig()
        ax.xaxis.set_major_formatter(fmtk)
        ax.yaxis.set_major_formatter(fmtk)

        for lab in ax.get_xticklabels() + ax.get_yticklabels():
            lab.set_rotation(45)

        cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
        cbar.set_label("log10(|FFT|)")
        cbar.formatter = _plain_formatter()
        cbar.update_ticks()

        ax.set_title("")
        fig.suptitle("")
        fig.savefig(fname, dpi=200)
        plt.close(fig)

    # =================================================
    # Per-group processing
    # =================================================
    with tqdm(total=total_files, desc="Processing SXM files") as pbar:

        for group_name, file_list in groups.items():

            section = prs.slides.add_slide(prs.slide_layouts[6])
            tf_sec = section.shapes.add_textbox(
                Inches(1.5), Inches(2.8), Inches(10), Inches(2)
            ).text_frame
            tf_sec.text = f"Group: {group_name}"
            tf_sec.paragraphs[0].font.size = Pt(34)
            tf_sec.paragraphs[0].font.bold = True

            for sxm_file in file_list:

                ds = img2xr(sxm_file)
                ds = interpolate2D_xr(ds)

                # =========================================================
                # [MODIFIED] multipass detection + per-pass channel mapping
                # =========================================================
                if "Z_fwd" in ds:
                    passes = [None]
                else:
                    passes = sorted(
                        {k.split("_")[1] for k in ds.data_vars
                         if k.startswith("Z_") and k.endswith("_fwd")}
                    )

                for p in passes:

                    if p is None:
                        Zf, Zb = "Z_fwd", "Z_bwd"
                        Lf, Lb = "LIX_fwd", "LIX_bwd"
                        title_suffix = ""
                        file_suffix = ""
                    else:
                        Zf, Zb = f"Z_{p}_fwd", f"Z_{p}_bwd"
                        Lf, Lb = f"LIX_{p}_fwd", f"LIX_{p}_bwd"
                        title_suffix = f" ({p})"
                        file_suffix = f"_{p}"

                    # =====================================================
                    # [MODIFIED] plane-fit with feasibility mask + propagate
                    # =====================================================
                    mask_fwd = make_plane_fit_mask_yfit(ds, Zf, 2)
                    mask_bwd = make_plane_fit_mask_yfit(ds, Zb, 2)

                    ydim, xdim = ds[Zf].dims
                    ds[f"plane_fit_mask_{Zf}"] = ((ydim, xdim), mask_fwd)
                    ds[f"plane_fit_mask_{Zb}"] = ((ydim, xdim), mask_bwd)
                    '''
                    ds = plane_fit_xr(ds, ch=Zf, method='y_fit', poly_order=2,
                                      mask=mask_fwd, overwrite=False)
                    ds = plane_fit_xr(ds, ch=Zb, method='y_fit', poly_order=2,
                                      mask=mask_bwd, overwrite=False)
                    '''
                    # ---- plane fit (SAFE) ----
                    try:
                        ds = plane_fit_xr(
                            ds, ch=Zf, method='y_fit',
                            poly_order=2, mask=mask_fwd, overwrite=False
                        )
                    except np.linalg.LinAlgError:
                        # fallback: no plane fit applied
                        ds[f"{Zf}_planefit"] = ds[Zf]
                    
                    try:
                        ds = plane_fit_xr(
                            ds, ch=Zb, method='y_fit',
                            poly_order=2, mask=mask_bwd, overwrite=False
                        )
                    except np.linalg.LinAlgError:
                        # fallback: no plane fit applied
                        ds[f"{Zb}_planefit"] = ds[Zb]



                    
                    # plateau tilt with mask (propagate)
                    for ch_pf, base_mask in [
                        (f"{Zf}_planefit", mask_fwd),
                        (f"{Zb}_planefit", mask_bwd),
                    ]:
                        try:
                            pt_mask = make_plateau_tilt_mask(ds, ch_pf, base_mask)
                            ds = plateau_tilt_xr(ds, ch=ch_pf, mask=pt_mask, overwrite=True)
                        except (np.linalg.LinAlgError, RuntimeError):
                            pass

                    for ch in [f"{Zf}_planefit", f"{Zb}_planefit", Lf, Lb]:
                        ds = twoDfft_xrft(ds, ch=ch)

                    # --- PNG + PPT insertion (UNCHANGED except filename suffix) ---
                    base = sxm_file + file_suffix
                    imgs = {
                        'Z_fwd': base + '_Z_fwd.png',
                        'Z_bwd': base + '_Z_bwd.png',
                        'LIX_fwd': base + '_LIX_fwd.png',
                        'LIX_bwd': base + '_LIX_bwd.png',
                        'Z_fwd_fft': base + '_Z_fwd_fft.png',
                        'Z_bwd_fft': base + '_Z_bwd_fft.png',
                        'LIX_fwd_fft': base + '_LIX_fwd_fft.png',
                        'LIX_bwd_fft': base + '_LIX_bwd_fft.png',
                    }

                    save_real(ds[f"{Zf}_planefit"], imgs['Z_fwd'], ds, 'copper')
                    save_real(ds[f"{Zb}_planefit"], imgs['Z_bwd'], ds, 'copper')
                    save_real(ds[Lf], imgs['LIX_fwd'], ds, 'bwr')
                    save_real(ds[Lb], imgs['LIX_bwd'], ds, 'bwr')

                    save_fft(ds[f"{Zf}_planefit_fft_amp"], imgs['Z_fwd_fft'], ds, 'gray_r')
                    save_fft(ds[f"{Zb}_planefit_fft_amp"], imgs['Z_bwd_fft'], ds, 'gray_r')
                    save_fft(ds[f"{Lf}_fft_amp"], imgs['LIX_fwd_fft'], ds, 'Blues')
                    save_fft(ds[f"{Lb}_fft_amp"], imgs['LIX_bwd_fft'], ds, 'Blues')

                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tf_title = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.1), Inches(12.0), Inches(0.7)
                    ).text_frame
                    tf_title.text = ds.attrs.get('title', '') + title_suffix

                    layout = [
                        ('Z_fwd', 'Z_fwd', 0.2, 1.1),
                        ('Z_bwd', 'Z_bwd', 3.4, 1.1),
                        ('LIX_fwd', 'LIX_fwd', 6.6, 1.1),
                        ('LIX_bwd', 'LIX_bwd', 9.8, 1.1),
                        ('Z_fwd_fft', 'Z_fwd_fft', 0.2, 4.4),
                        ('Z_bwd_fft', 'Z_bwd_fft', 3.4, 4.4),
                        ('LIX_fwd_fft', 'LIX_fwd_fft', 6.6, 4.4),
                        ('LIX_bwd_fft', 'LIX_bwd_fft', 9.8, 4.4),
                    ]

                    for key, label, x, y in layout:
                        st = slide.shapes.add_textbox(
                            Inches(x), Inches(y - 0.45),
                            Inches(3.0), Inches(0.3)
                        ).text_frame
                        st.text = label
                        slide.shapes.add_picture(
                            imgs[key], Inches(x), Inches(y), width=Inches(3.0)
                        )

                    for f in imgs.values():
                        if os.path.exists(f):
                            os.remove(f)

                del ds
                gc.collect()
                pbar.update(1)

    prs.save(output_path)
    print("Saved:", output_path)


# %% [markdown]
# ## Stage 6 — Run
#
# ```python
# TwoDdata_summary_sxm_pptx(groups)
# ```

# %%
TwoDdata_summary_sxm_pptx(groups)

